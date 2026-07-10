import io
from typing import List

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text_from_bytes(file_bytes: bytes, filename: str) -> str:
    lower = filename.lower()
    if lower.endswith('.pdf'):
        text = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)
    if lower.endswith('.docx'):
        doc = DocxDocument(io.BytesIO(file_bytes))
        return '\n'.join(p.text for p in doc.paragraphs)
    if lower.endswith('.pptx'):
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
        return '\n'.join(texts)
    return file_bytes.decode('utf-8', errors='ignore')


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    tokens = text.split()
    chunks: List[str] = []
    i = 0
    while i < len(tokens):
        chunk = tokens[i:i+chunk_size]
        chunks.append(' '.join(chunk))
        i += chunk_size - overlap
    return chunks
